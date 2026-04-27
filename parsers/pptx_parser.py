import logging
from pathlib import Path

from pptx import Presentation

from parsers.base import ParsedSegment

log = logging.getLogger(__name__)


def parse_pptx(path: str | Path) -> list[ParsedSegment]:
    """Извлекает текст по слайдам, сохраняя номер слайда в метаданных.
    Картинки игнорируются."""
    segments: list[ParsedSegment] = []
    try:
        prs = Presentation(str(path))
    except Exception as e:
        log.error("Не удалось открыть pptx %s: %s", path, e)
        return []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"\n[Заметки докладчика]\n{notes}")

        text = "\n".join(parts).strip()
        if text:
            segments.append(
                ParsedSegment(
                    text=text,
                    slide_number=i,
                    metadata={"source_type": "pptx"},
                )
            )
    return segments
